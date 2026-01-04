from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def create_pitch_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    primary_blue = RGBColor(30, 64, 175)
    accent_green = RGBColor(16, 185, 129)
    dark_text = RGBColor(31, 41, 55)
    
    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = primary_blue
        background.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 220, 255)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, bullet_points):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = primary_blue
        
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = accent_green
        line.line.fill.background()
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(22)
            p.font.color.rgb = dark_text
            p.space_before = Pt(12)
            p.space_after = Pt(8)
        
        return slide
    
    def add_two_column_slide(title, left_title, left_points, right_title, right_points):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = primary_blue
        
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = accent_green
        line.line.fill.background()
        
        left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.6))
        tf = left_header.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = accent_green
        
        left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(4.5))
        tf = left_content.text_frame
        tf.word_wrap = True
        for i, point in enumerate(left_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = dark_text
            p.space_before = Pt(8)
        
        right_header = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.6))
        tf = right_header.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = accent_green
        
        right_content = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(5.5), Inches(4.5))
        tf = right_content.text_frame
        tf.word_wrap = True
        for i, point in enumerate(right_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = dark_text
            p.space_before = Pt(8)
        
        return slide
    
    add_title_slide(
        "Logos AI",
        "AI Phone System for SMB Retail & eCommerce\nReducing support costs by automating repetitive inbound calls"
    )
    
    add_content_slide("Mission", [
        "Logos AI helps SMB retail and eCommerce businesses cut support costs by automating repetitive inbound phone calls with AI",
        "We focus on what enterprises overlook: the 10-500 employee businesses drowning in simple calls",
        "Our AI answers order status, pickup readiness, store hours & FAQs - and knows when to hand off to humans"
    ])
    
    add_content_slide("The Founder", [
        "Fuzail Kadri - Founder & CEO",
        "Mechanical Design Engineer by education with 4+ years in SaaS",
        "Customer Success Manager with deep experience in Customer Support and Contact Centres",
        "Combined commercial SaaS expertise with passion for automation (built open-source service robot as engineering project)",
        "Understands customer pain points from both sides of the support equation"
    ])
    
    add_content_slide("The Problem", [
        "SMBs spend too much time & money answering repetitive inbound calls (order status, store hours, basic policy questions)",
        "Calls interrupt in-store staff serving customers, reducing sales effectiveness",
        "Calls go to voicemail during peak hours or after-hours, frustrating customers",
        "$100K+ lost per location annually from missed calls (BIA/Kelsey research)",
        "Current solutions (IVRs, chatbots, FAQs) don't address PHONE calls effectively"
    ])
    
    add_two_column_slide(
        "Current Solutions Fall Short",
        "What SMBs Do Today",
        [
            "Manually answer every call with staff",
            "Use basic IVRs that frustrate callers",
            "Add FAQs/chatbots that don't help phone callers",
            "Hire additional staff or outsource",
            "Let calls go to voicemail"
        ],
        "The Result",
        [
            "Staff interrupted from serving customers",
            "Increased costs without solving the problem",
            "Missed revenue opportunities",
            "Poor customer experience",
            "Teams stuck on low-value work"
        ]
    )
    
    add_content_slide("Our Solution", [
        "AI system that answers business phone calls for common questions",
        "2 AI intents live today: order status (via CSV) and store hours",
        "Intelligently escalates complex issues to human agents",
        "Simple setup: upload your order data via CSV, customize greetings, go live",
        "Integrates with Twilio for reliable telephony",
        "Listen Mode analytics and voicemail recording included"
    ])
    
    add_content_slide("Why Logos AI is Different", [
        "Built for SMBs, not enterprises - pricing and features designed for 10-500 employee businesses",
        "Phone-first approach - we solve the phone call problem, not another chatbot",
        "Simple setup - CSV upload, not complex API integrations required",
        "Smart handoff - AI knows its limits and transfers to humans when needed",
        "Transparent pricing - subscription model with no hidden fees"
    ])
    
    add_content_slide("Business Model", [
        "Subscription-based SaaS model",
        "Monthly fee based on call volume",
        "Simple plans designed for SMB budgets",
        "No hidden fees or surprise charges",
        "Optional annual discounts for committed customers",
        "Pilot program at $99/month for early adopters"
    ])
    
    add_content_slide("Market Opportunity", [
        "AI Agents Market: $7.63B (2025) → $50.31B by 2030 (35-40% CAGR)",
        "Voice AI for Restaurants alone: $10B (2025) → $49B by 2029",
        "SMB retail represents massive untapped opportunity",
        "Competitors focus on enterprise; SMBs are underserved",
        "90% of hospitals expected to use AI agents by end of 2025",
        "Clear expansion path to adjacent verticals"
    ])
    
    add_two_column_slide(
        "Expansion Roadmap",
        "Phase 1: Retail & eCommerce (Current)",
        [
            "Order status inquiries",
            "Store hours & location info",
            "Pickup readiness checks",
            "Basic policy FAQs",
            "Intelligent human handoff"
        ],
        "Phase 2: Q1-Q2 2025",
        [
            "Shopify integration",
            "Returns & exchanges automation",
            "Stripe payment integration",
            "Subscription management",
            "Multi-location support"
        ]
    )
    
    add_two_column_slide(
        "Industry Expansion (2025-2026)",
        "Restaurants & Hospitality",
        [
            "Reservation management",
            "Order taking (phone orders)",
            "Menu inquiries & specials",
            "Wait time updates",
            "$100B+ lost annually to missed restaurant calls"
        ],
        "Healthcare & Services",
        [
            "Appointment scheduling",
            "Insurance verification",
            "Prescription refill status",
            "Lab result notifications",
            "Healthcare AI: $38.66B market (2025)"
        ]
    )
    
    add_content_slide("Traction & Proof Points", [
        "Live working prototype ready for pilot customers",
        "2 AI intents working today: order status (CSV), store hours",
        "Twilio telephony integration operational",
        "Human handoff + voicemail + Listen Mode analytics all functional",
        "FAQ responses planned for Q1 2025",
        "Shopify integration planned for Q1 2025"
    ])
    
    add_two_column_slide(
        "Financial Projections",
        "Year 1 Targets",
        [
            "10-20 pilot customers at $99/mo",
            "ARR target: $12,000 - $24,000",
            "Focus on product-market fit",
            "Prove unit economics",
            "Build case studies"
        ],
        "Year 2-3 Growth",
        [
            "100+ customers across retail/restaurants",
            "ARR target: $300,000 - $500,000",
            "Launch Shopify & restaurant integrations",
            "Expand to hospitality vertical",
            "Achieve profitability path"
        ]
    )
    
    add_content_slide("Team & Advisors", [
        "Fuzail Kadri - Founder & CEO",
        "   - Mechanical Design Engineer with 4+ years SaaS experience",
        "   - Customer Success Manager background - understands customer pain points",
        "   - Built open-source service robot as engineering capstone",
        "",
        "Seeking advisors with expertise in:",
        "   - AI/ML and voice technology",
        "   - SMB SaaS go-to-market",
        "   - Retail and eCommerce operations"
    ])
    
    add_content_slide("Go-to-Market Strategy", [
        "Phase 1: Direct outreach to SMB retail/eCommerce (Shopify stores, local retailers)",
        "Phase 2: Partner with Shopify app ecosystem and POS providers",
        "Phase 3: Expand to restaurant aggregators and hospitality networks",
        "Customer acquisition through:",
        "   - Content marketing (call handling ROI calculators, case studies)",
        "   - Industry-specific trade shows and events",
        "   - Referral program for satisfied pilot customers"
    ])
    
    add_two_column_slide(
        "Use of Funds",
        "Product Development (60%)",
        [
            "Shopify & eCommerce integrations",
            "Restaurant industry features",
            "Advanced AI training & improvements",
            "Multi-language support",
            "Mobile app for business owners"
        ],
        "Go-to-Market (40%)",
        [
            "Sales & marketing team",
            "Customer success resources",
            "Pilot program expansion",
            "Trade show presence",
            "Content & SEO investment"
        ]
    )
    
    add_content_slide("The Ask", [
        "Seeking $500K seed funding to accelerate growth",
        "",
        "Milestones this funding enables:",
        "   - 50+ paying customers within 12 months",
        "   - Shopify integration launch (Q1 2025)",
        "   - Restaurant vertical expansion (Q2 2025)",
        "   - Proven unit economics and path to Series A"
    ])
    
    add_title_slide(
        "Thank You",
        "Fuzail Kadri | Logos AI\nfuzail@logosai.com"
    )
    
    os.makedirs('static/downloads', exist_ok=True)
    output_path = 'static/downloads/logos_ai_pitch_deck.pptx'
    prs.save(output_path)
    
    return output_path

if __name__ == "__main__":
    path = create_pitch_deck()
    print(f"Pitch deck created: {path}")
